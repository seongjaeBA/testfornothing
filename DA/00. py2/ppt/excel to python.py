import os
import pptx
from pptx import Presentation

## 슬라이드 마스터 index
SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_TITLE_AND_CONTENT = 1
SLD_LAYOUT_Section_Header = 2
SLD_LAYOUT_Two_Content = 3
SLD_LAYOUT_Comparison = 4
SLD_LAYOUT_Title_only = 5
SLD_LAYOUT_Blank = 6
SLD_LAYOUT_Content_Caption = 7
SLD_LAYOUT_Picture_Caption = 8

pic_dir = './picture'

picture_list = os.listdir(pic_dir)


#파일 불러오기
'''f = open('./test.pptx')'''
prs = Presentation('./test.pptx')

# or
'''
with open('foobar.pptx') as f:
    source_stream = StringIO(f.read())
prs = Presentation(source_stream)
source_stream.close()
...
target_stream = StringIO()
prs.save(target_stream)
'''

for i in picture_list:
    #슬라이드 레이아웃 set
    slide_layout = prs.slide_layouts[8]
    #슬라이드 더하기
    slide = prs.slides.add_slide(slide_layout)

    '''
    #shape 리스트에 접근
    shapes = slide.shapes
    left = top = width = height = Inches(1.0)
    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        ) 
    #shape 채우기
    shape.fill.background()
    '''
    #place-holder 알아보기
    '''
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    '''

    ##place-holder 접근
    #slide.placeholder[idx] - name 

    ## 그림 placeholder를 찾아서 접근.
    picture = slide.placeholders[1].insert_picture('./picture/%s' % (i))
    


#파일로 저장(파일로 하기)
prs.save('./test1.pptx')
'''f.close'''
