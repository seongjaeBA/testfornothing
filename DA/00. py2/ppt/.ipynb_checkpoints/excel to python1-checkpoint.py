import os
import csv
import pptx
import pandas as pd
from pptx import Presentation
from pptx import slide


#파일 불러오기
file_list =pd.read_csv('c:/', names=[], header=[], index_col=)
#file_list = pd.read_excel(''. sheetname=)

pic_dir = './picture'
picture_list = os.listdir(pic_dir)
#picture_list = df[].tolist()
pptx_dir = './ppt'
pptx_list = os.listdir(pptx_dir)
for pt in pptx_list:
    prs = Presentation('./ppt/%s' %(pt))

    slides = prs.slides
    print(len(slides))

    for i in picture_list:
        for idx in range(len(slides)):
            #print(slides[idx].shapes)
            slide = slides[idx]
            for shape in slide.placeholders:
                print('%d %s' % (shape.placeholder_format.idx, shape.name))
                picture = slide.placeholders[1].insert_picture('./picture/%s' % (i))


#파일로 저장(새파일로 하기)
    prs.save('PI_'+pt)

